"""python-pptx based PPT parser. Extracts text frames + tables + notes (no pictures)."""
from io import BytesIO

from pptx import Presentation


def parse_pptx(buffer: bytes) -> dict:
    prs = Presentation(BytesIO(buffer))
    slides = []
    raw_parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []
        notes = ""
        if slide.shapes.title:
            title = slide.shapes.title.text or ""
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                body_parts.append(shape.text_frame.text)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        body_parts.append(cell.text_frame.text)
            # Skip picture elements; embedded-image OCR is intentionally unsupported.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text or ""
        body = "\n".join(body_parts).strip()
        slides.append({"index": idx, "title": title, "body": body, "notes": notes})
        raw_parts.append(f"--- SLIDE {idx} ---\n# {title}\n\n{body}\n\nNOTES: {notes}")
    return {
        "slides": slides,
        "slide_count": len(slides),
        "raw_text": "\n\n".join(raw_parts),
    }
